import os
import json
import re
from langchain.prompts import PromptTemplate
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import PyPDFLoader
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_core.prompts import ChatPromptTemplate
from langchain.docstore.document import Document as LangchainDocument
from docx import Document as DocxDocument
from pptx import Presentation

from app.core.settings import settings

# Initialize the Gemini LLM
llm = ChatGoogleGenerativeAI(
    model="gemini-1.5-pro",
    api_key=settings.gemini_key,
    temperature=0.5,
    max_tokens=1024,
    timeout=30,
    max_retries=3,
)

# Define prompt templates
question_prompt = PromptTemplate(
    input_variables=["content", "num_questions", "question_type"],
    template=(
        "Generate {num_questions} {question_type} questions based on the following content:\n"
        "\n{content}\n\nFormat: Question text, options (if any), and correct answer."
    ),
)

summarization_prompt = ChatPromptTemplate.from_messages(
    [
        ("system", "You are an expert summarizer. Summarize the following text."),
        ("human", "{content}"),
    ]
)


def load_docx(file_path):
    """Loads text from a .docx file."""
    doc = DocxDocument(file_path)
    content = "\n".join([paragraph.text for paragraph in doc.paragraphs])
    return [LangchainDocument(page_content=content)]


def load_pptx(file_path):
    """Loads text from a .pptx file."""
    presentation = Presentation(file_path)
    content = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                content.append(shape.text)
    full_content = "\n".join(content)
    return [LangchainDocument(page_content=full_content)]


# Loaders for different file types
FILE_LOADERS = {
    ".pdf": PyPDFLoader,
    ".docx": load_docx,
    ".pptx": load_pptx,
}


def extract_text_from_file(file_path):
    """Extracts text from various file formats."""
    file_extension = os.path.splitext(file_path)[1].lower()
    if file_extension in FILE_LOADERS:
        loader = FILE_LOADERS[file_extension]
        if file_extension == ".pdf":
            documents = loader(file_path).load()
        else:
            documents = loader(file_path)
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000, chunk_overlap=100
        )
        return text_splitter.split_documents(documents)
    else:
        raise ValueError(f"Unsupported file format: {file_extension}")


def generate_questions(content, num_questions, question_type):
    """Generates questions using LangChain and Gemini API."""
    # Create the summarization chain
    summarization_chain = summarization_prompt | llm

    # Summarize the content
    summarized_content = summarization_chain.invoke({"content": content})

    # Extract summarized text
    summarized_text = (
        summarized_content.content
        if hasattr(summarized_content, "content")
        else str(summarized_content)
    )

    # Define the question generation chain
    question_chain = question_prompt | llm

    # Generate questions
    questions = question_chain.invoke(
        {
            "content": summarized_text,
            "num_questions": num_questions,
            "question_type": question_type,
        }
    )

    # Return the generated questions
    return questions.content if hasattr(questions, "content") else str(questions)


def format_questions_to_json(questions_content):
    """Converts the questions content into JSON format."""
    questions_list = []

    # Split the questions content into individual questions (assuming each question is separated by two newlines)
    question_blocks = questions_content.split("\n\n")

    # Parse each question block
    for question_block in question_blocks:
        # Extract the question text
        question_match = re.match(r"(.*?)(?=\n[A-D])", question_block.strip())
        if question_match:
            question_text = question_match.group(1).strip()
        else:
            continue

        # Extract the choices (assuming they are labeled as A, B, C, D, etc.)
        choices = []
        for letter in ["A", "B", "C", "D"]:
            choice_match = re.search(rf"{letter}\.\s*(.*?)\n", question_block)
            if choice_match:
                choices.append(choice_match.group(1).strip())

        # Extract the correct answer (assuming the correct answer is labeled like "Correct Answer: B")
        correct_answer_match = re.search(r"Correct Answer:\s*(\w)", question_block)
        if correct_answer_match:
            correct_answer = ord(correct_answer_match.group(1).upper()) - ord("A")
        else:
            correct_answer = 0  # Default to the first choice if no answer is provided

        # Build the question data
        question_data = {
            "question": question_text,
            "choices": choices,
            "correctAnswer": correct_answer,
        }

        questions_list.append(question_data)

    # Return the JSON structure as a string
    return json.dumps(questions_list, indent=4)


def save_output(output, output_format, file_name):
    """Saves the generated questions in the desired format."""
    if isinstance(output, str):  # Validate output is a string
        if output_format == "markdown":
            with open(file_name, "w") as f:
                f.write(output)
        elif output_format == "html":
            html_content = f"<html><body><pre>{output}</pre></body></html>"
            with open(file_name, "w") as f:
                f.write(html_content)
        elif output_format == "json":
            kahoot_json = format_questions_to_json(output)
            with open(file_name, "w") as f:
                json.dump(kahoot_json, f, indent=4)
        else:
            raise ValueError("Unsupported output format.")
    else:
        raise TypeError("Output must be a string.")
